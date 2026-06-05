"""
Tools Engine
Provides local capabilities (superpowers) for the AI, like running Python code.
"""
import subprocess
import tempfile
import os
import logging
from typing import Dict, Any

logger = logging.getLogger(__name__)

class ToolResponse:
    """Standardized result wrapper for all neural tools."""
    def __init__(self, success: bool, output: str, metadata: dict = None):
        self.success = success
        self.output = output
        self.metadata = metadata or {}

class ToolRegistry:
    """Governance layer for tool permissions and execution."""
    PERMISSIONS = {
        "execute_python_code": "EXECUTE",
        "create_pptx": "WRITE",
        "web_scrape": "READ",
        "read_file": "READ",
        "write_file": "WRITE",
        "calculate": "READ",
        "get_datetime": "READ"
    }
    
    @classmethod
    def validate_permission(cls, tool_name: str, requested_by: str = "JARVIS") -> bool:
        """Autonomous check for tool execution authority."""
        permission = cls.PERMISSIONS.get(tool_name)
        if not permission:
            logger.warning(f"Governance Breach: Tool '{tool_name}' not registered in PERMISSIONS grid.")
            return False
            
        logger.info(f"[{requested_by}] Authorization Granted: {tool_name} ({permission})")
        return True

class SecurityGuard:
    """Neural Firewall: Scrubs secrets from logs and outputs."""
    BANNED_PATTERNS = [
        r"sk-[a-zA-Z0-9]{48}", # OpenAI
        r"gsk_[a-zA-Z0-9]{48}", # Groq
        r"pk_test_[a-zA-Z0-9]{24}", # Stripe
        r"sk_test_[a-zA-Z0-9]{24}"  # Stripe
    ]
    
    @classmethod
    def scrub(cls, text: str) -> str:
        """Anonymize sensitive tokens in string data."""
        import re
        scrubbed = text
        for pattern in cls.BANNED_PATTERNS:
            scrubbed = re.sub(pattern, "[REDACTED_SECRET]", scrubbed)
        return scrubbed

class ToolEngine:
    @staticmethod
    def create_pptx(slides_content: list, filename: str = "presentation.pptx") -> str:
        """
        Creates a PowerPoint file from a list of slide dicts.
        Each dict should have 'title', 'content' (list), and 'speaker_notes'.
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            
            prs = Presentation()
            
            # Slide layouts: 0 is Title, 1 is Title and Content
            for slide_data in slides_content:
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                
                title = slide.shapes.title
                content = slide.placeholders[1]
                
                title.text = slide_data.get("title", "Untitled Slide")
                
                body = content.text_frame
                content_items = slide_data.get("content", [])
                if isinstance(content_items, str):
                    body.text = content_items
                else:
                    for item in content_items:
                        p = body.add_paragraph()
                        p.text = str(item)
                        p.level = 0
                
                # Speaker Notes
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_data.get("speaker_notes", "")

            # Save to a generic exports directory
            export_dir = os.path.join(os.getcwd(), "data", "exports")
            os.makedirs(export_dir, exist_ok=True)
            
            # Ensure unique filename
            from datetime import datetime
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            safe_filename = f"jarvis_pres_{timestamp}.pptx"
            save_path = os.path.join(export_dir, safe_filename)
            
            prs.save(save_path)
            logger.info(SecurityGuard.scrub(f"PowerPoint saved: {save_path}"))
            return safe_filename
            
        except Exception as e:
            logger.error(SecurityGuard.scrub(f"Error creating PPTX: {e}"))
            return f"Error: {e}"

    @staticmethod
    def execute_python_code(code: str, timeout: int = 15) -> str:
        """
        Executes Python code locally with structural AST security analysis.
        """
        if not ToolRegistry.validate_permission("execute_python_code"):
            return "Error: Governance Access Denied."
            
        import ast
        logger.info("Neural Guard: Performing AST Structural Analysis...")
        
        # ─── AST Security Firewall ──────────────────────────────────────────
        try:
            tree = ast.parse(code)
            
            # Dangerous nodes to block
            BLOCKED_NODES = (ast.Import, ast.ImportFrom) 
            DANGEROUS_FUNCTIONS = {"eval", "exec", "open", "input", "__import__", "getattr", "setattr"}
            DANGEROUS_MODULES = {"os", "subprocess", "socket", "shutil", "requests", "urllib", "sys", "pty"}

            for node in ast.walk(tree):
                # 1. Block all imports (user should not import system modules)
                if isinstance(node, BLOCKED_NODES):
                    return "Security Breach: Arbitrary imports are restricted in the neural sandbox."
                
                # 2. Block calls to dangerous built-ins
                if isinstance(node, ast.Call):
                    if isinstance(node.func, ast.Name) and node.func.id in DANGEROUS_FUNCTIONS:
                        return f"Security Breach: Execution of restricted function '{node.func.id}' aborted."
                    
                    # Block calls to dangerous attributes (e.g., os.system)
                    if isinstance(node.func, ast.Attribute):
                        attr_name = node.func.attr
                        if attr_name in ["system", "popen", "spawn", "call", "run"]:
                            return f"Security Breach: Execution of restricted attribute '{attr_name}' aborted."
                
                # 3. Block accessing internal or magic names (__*)
                if isinstance(node, ast.Name) and node.id.startswith("__"):
                    return f"Security Breach: Access to internal identifier '{node.id}' aborted."
                if isinstance(node, ast.Attribute) and node.attr.startswith("__"):
                    return f"Security Breach: Access to internal attribute '{node.attr}' aborted."

            logger.info("AST Analysis: CLEAN. Proceeding to execution.")
            
        except SyntaxError as e:
            return f"Syntax Error in neural logic: {e}"
        except Exception as e:
            return f"Security Guard Error: {e}"

        # ─── Execution Logic ────────────────────────────────────────────────

        # Create a temporary file to hold the code
        with tempfile.NamedTemporaryFile(mode='w', suffix='.py', delete=False) as temp_file:
            # Add basic safe imports if they aren't there but don't mess with user's structure too much
            temp_file.write(code)
            temp_path = temp_file.name

        try:
            # Run the code
            result = subprocess.run(
                ["python", temp_path],
                capture_output=True,
                text=True,
                timeout=timeout,
                check=False
            )
            
            output = ""
            if result.stdout:
                output += "--- STDOUT ---\n" + result.stdout + "\n"
            if result.stderr:
                output += "--- STDERR ---\n" + result.stderr + "\n"
            
            if not output:
                output = "Code executed successfully with no output."
                
            return output
            
        except subprocess.TimeoutExpired:
            return f"Error: Code execution timed out after {timeout} seconds."
        except Exception as e:
            return f"Error executing code: {e}"
        finally:
            # Clean up the temp file
            if os.path.exists(temp_path):
                try:
                    os.remove(temp_path)
                except:
                    pass

    # ─── Web Scraping ───────────────────────────────────────────────

    @staticmethod
    def web_scrape(url: str) -> str:
        """Scrape main content from a URL."""
        try:
            import requests
            from bs4 import BeautifulSoup

            headers = {
                "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
                              "AppleWebKit/537.36 Chrome/131.0 Safari/537.36"
            }
            resp = requests.get(url, headers=headers, timeout=15)
            resp.raise_for_status()

            soup = BeautifulSoup(resp.text, "lxml")

            # Remove unwanted elements
            for tag in soup.find_all(["script", "style", "nav", "footer", "aside", "iframe", "header"]):
                tag.decompose()

            # Try to find main content
            main = None
            for selector in ["article", "main", '[role="main"]', ".post-content", ".article-body"]:
                main = soup.select_one(selector)
                if main:
                    break

            if main:
                text = main.get_text(separator="\n", strip=True)
            else:
                paragraphs = soup.find_all("p")
                text = "\n\n".join(
                    p.get_text(strip=True) for p in paragraphs
                    if len(p.get_text(strip=True)) > 30
                )

            # Clean up
            import re
            text = re.sub(r'\n{3,}', '\n\n', text)
            text = re.sub(r' {2,}', ' ', text)

            if len(text) > 3000:
                text = text[:3000] + "\n\n[... content truncated]"

            return text if text else "Could not extract meaningful content from this URL."

        except Exception as e:
            logger.error(f"Scrape error for {url}: {e}")
            return f"Error scraping URL: {e}"

    # ─── File Operations ────────────────────────────────────────────

    @staticmethod
    def read_file(filepath: str) -> str:
        """Read contents of a file."""
        try:
            # Security: only allow reading from data directory
            import config as cfg
            allowed_dir = os.path.abspath(cfg.DATA_DIR)
            abs_path = os.path.abspath(filepath)

            # Allow relative paths from data dir
            if not abs_path.startswith(allowed_dir):
                # Try with data dir prefix
                abs_path = os.path.join(allowed_dir, filepath)
                if not os.path.abspath(abs_path).startswith(allowed_dir):
                    return "Error: File access restricted to the data directory."

            if not os.path.exists(abs_path):
                return f"File not found: {filepath}"

            with open(abs_path, "r", encoding="utf-8") as f:
                content = f.read()

            if len(content) > 5000:
                content = content[:5000] + "\n\n[... file truncated]"

            return content

        except Exception as e:
            return f"Error reading file: {e}"

    @staticmethod
    def write_file(filepath: str, content: str) -> str:
        """Write content to a file in the data directory."""
        try:
            import config as cfg
            export_dir = os.path.join(cfg.DATA_DIR, "exports")
            os.makedirs(export_dir, exist_ok=True)

            # Force all writes into the exports directory
            safe_name = os.path.basename(filepath)
            save_path = os.path.join(export_dir, safe_name)

            with open(save_path, "w", encoding="utf-8") as f:
                f.write(content)

            logger.info(f"File written: {save_path}")
            return f"File saved successfully: {safe_name}"

        except Exception as e:
            return f"Error writing file: {e}"

    # ─── Calculator ─────────────────────────────────────────────────

    @staticmethod
    def calculate(expression: str) -> str:
        """Safely evaluate a math expression."""
        try:
            import ast
            import operator
            import math

            # Supported operators
            operators = {
                ast.Add: operator.add,
                ast.Sub: operator.sub,
                ast.Mult: operator.mul,
                ast.Div: operator.truediv,
                ast.Pow: operator.pow,
                ast.Mod: operator.mod,
                ast.USub: operator.neg,
                ast.UAdd: operator.pos,
            }
            
            # Supported functions
            functions = {
                "abs": abs, "round": round, "min": min, "max": max,
                "pow": pow, "int": int, "float": float,
                "sqrt": math.sqrt, "sin": math.sin, "cos": math.cos, 
                "tan": math.tan, "log": math.log, "log10": math.log10,
            }
            
            # Supported constants
            constants = {
                "pi": math.pi,
                "e": math.e
            }

            def evaluate(node):
                if isinstance(node, ast.Constant):
                    return node.value
                elif isinstance(node, ast.BinOp):
                    op = type(node.op)
                    if op not in operators:
                        raise ValueError(f"Unsupported operator: {op}")
                    return operators[op](evaluate(node.left), evaluate(node.right))
                elif isinstance(node, ast.UnaryOp):
                    op = type(node.op)
                    if op not in operators:
                        raise ValueError(f"Unsupported operator: {op}")
                    return operators[op](evaluate(node.operand))
                elif isinstance(node, ast.Call):
                    if not isinstance(node.func, ast.Name) or node.func.id not in functions:
                        raise ValueError(f"Unsupported function: {node.func.id if isinstance(node.func, ast.Name) else 'unknown'}")
                    args = [evaluate(arg) for arg in node.args]
                    return functions[node.func.id](*args)
                elif isinstance(node, ast.Name):
                    if node.id in constants:
                        return constants[node.id]
                    raise ValueError(f"Unsupported variable: {node.id}")
                else:
                    raise TypeError(f"Unsupported expression node: {type(node)}")

            # Parse and evaluate
            tree = ast.parse(expression, mode='eval')
            result = evaluate(tree.body)
            
            return f"{expression} = {result}"

        except Exception as e:
            return f"Calculation error: {e}"

    # ─── Date/Time ──────────────────────────────────────────────────

    @staticmethod
    def get_datetime() -> str:
        """Get current date and time."""
        from datetime import datetime, timezone
        now = datetime.now()
        utc = datetime.now(timezone.utc)
        return (
            f"Local time: {now.strftime('%A, %B %d, %Y at %I:%M %p')}\n"
            f"UTC time: {utc.strftime('%Y-%m-%d %H:%M:%S UTC')}"
        )

    # ─── Tool Registry ──────────────────────────────────────────────

    @classmethod
    def list_tools(cls) -> list:
        """List all available tools."""
        return [
            {"name": "execute_python_code", "description": "Run Python code locally"},
            {"name": "create_pptx", "description": "Create PowerPoint presentations"},
            {"name": "web_scrape", "description": "Extract content from a URL"},
            {"name": "read_file", "description": "Read a file from data directory"},
            {"name": "write_file", "description": "Write/save a file to data directory"},
            {"name": "calculate", "description": "Evaluate math expressions"},
            {"name": "get_datetime", "description": "Get current date and time"},
        ]
